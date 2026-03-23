#!/usr/bin/env python3
"""培训课件生成器 - 两阶段生成（大纲→内容），26.67×15 英寸大画布"""
import os, sys, json, requests, time, argparse
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

API_KEY  = os.environ.get("AIGC_API_KEY", "")
API_BASE = os.environ.get("AIGC_BASE_URL", "https://aigc.sankuai.com/v1/openai/native")
MODEL    = os.environ.get("AIGC_MODEL", "gpt-4.1")
HEADERS  = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}

# ── 配色 ─────────────────────────────────────────────
C = {
    "accent":   RGBColor(0x1F, 0x6F, 0xB8),
    "accent2":  RGBColor(0x2E, 0xAA, 0x76),
    "dark":     RGBColor(0x1A, 0x1A, 0x2E),
    "gray":     RGBColor(0x55, 0x55, 0x6B),
    "lgray":    RGBColor(0xF4, 0xF5, 0xF7),
    "white":    RGBColor(0xFF, 0xFF, 0xFF),
    "line":     RGBColor(0xD0, 0xD5, 0xDD),
    "hilite":   RGBColor(0xFF, 0xF3, 0xCC),
    "green_bg": RGBColor(0xE8, 0xF5, 0xE9),
    "red_bg":   RGBColor(0xFD, 0xED, 0xED),
}

# ── API ──────────────────────────────────────────────
def chat(messages, temperature=0.5):
    payload = {"model": MODEL, "messages": messages, "temperature": temperature, "max_tokens": 6000}
    r = requests.post(f"{API_BASE}/chat/completions", headers=HEADERS, json=payload, timeout=180)
    r.raise_for_status()
    return r.json()["choices"][0]["message"]["content"].strip()

def parse_json(raw):
    raw = raw.strip()
    for fence in ("```json", "```"):
        if raw.startswith(fence):
            raw = raw[len(fence):]
            if raw.endswith("```"): raw = raw[:-3]
            break
    return json.loads(raw.strip())

# ════════════════════════════════════════════════════
# 阶段一：生成课程大纲
# ════════════════════════════════════════════════════
OUTLINE_SYSTEM = """你是专业培训课程设计师，擅长将用户需求转化为结构化的课程大纲。

## 场景类型
- cover: 封面页
- section: 章节过渡页（每章开头）
- content: 知识讲解页
- example: 实战案例/场景演练页
- table: 表格对比/流程规范页
- quiz: 互动问答/知识检验页
- summary: 总结复盘页

## 设计原则
1. 每个页面有明确的教学目标
2. 逻辑递进：从概念→原则→案例→练习→总结
3. 站在学员视角设计学习体验
4. content页：4~6个要点，每个要点20~40字，配100~150字详细说明
5. example页：必须包含完整对话场景，含❌错误示范和✅正确示范，共150~200字
6. table页：填充5~8行完整数据
7. summary页：5~6条核心结论，每条15~25字

## 输出格式
JSON数组，每项：
{
  "id": "scene_1",
  "type": "cover|section|content|example|table|quiz|summary",
  "title": "页面标题",
  "subtitle": "副标题（cover/section页）",
  "teaching_objective": "本页教学目标（一句话）",
  "key_points": ["要点1（20-40字）", "要点2", ...],
  "detail": "详细说明（content页，100-150字，解释背景/原因/操作步骤）",
  "scenario": "对话场景（example页，含❌错误示范和✅正确示范，共150-200字）",
  "table_headers": ["列名1","列名2","列名3"],
  "table_rows": [["值1","值2","值3"],...],
  "quiz_question": "问题（quiz页）",
  "quiz_options": ["A. 选项1","B. 选项2","C. 选项3","D. 选项4"],
  "quiz_answer": "A",
  "quiz_explain": "解析（50字内）",
  "layout_hint": "布局建议：如'左栏要点+右栏说明'/'全页场景对话'/'三列对比卡片'"
}

只返回JSON数组。"""

OUTLINE_USER = """请为「{topic}」生成 {n_slides} 页培训PPT大纲。

页面分配建议（{n_slides}页）：
- 第1页：cover（封面）
- 第2页：section（第一章引导，key_points填本章内容预告4-5条）
- 第3-5页：content（核心知识，每页4-5要点+完整detail说明）
- 第6页：table（规范/流程/对比表格）
- 第7-8页：example（实战案例，含完整对话示范）
- 第9页：section（第二章引导）
- 第10-11页：content（知识讲解）
- 第12页：quiz（互动问答）
- 第13页：summary（总结）

如页数不同请按比例调整，保持结构完整。
只返回JSON数组，不要其他内容。"""

# ════════════════════════════════════════════════════
# 阶段二：为每页生成详细内容
# ════════════════════════════════════════════════════
ENRICH_SYSTEM = """你是资深培训内容专家。给你一个PPT页面的大纲信息，你需要将其扩充为完整、饱满的培训内容。

## 内容质量要求
- content页的每个key_point必须具体、可操作，包含实际场景或数据
- detail要涵盖：是什么 → 为什么 → 怎么做，有逻辑层次
- scenario要真实还原客服对话场景，❌示范要指出具体问题，✅示范要给出完整规范话术
- table数据要完整、实用，不要填占位符
- quiz选项要有一定迷惑性，answer_explain要点出考察点

## 输出
在原始大纲JSON基础上，只补充/扩充内容字段，保持原有结构。
返回单个JSON对象（不是数组）。"""

ENRICH_USER = """请扩充以下PPT页面内容，使其更饱满实用：

主题：{topic}
页面大纲：
{outline_item}

要求：
- 如果是content页：确保每个key_point 20-40字，detail 120-150字
- 如果是example页：scenario必须包含完整对话（❌错误示范+✅正确示范），总计150-200字
- 如果是table页：确保table_rows有5-8行完整数据
- 如果是section页：key_points填充本章内容预告，4-5条
- 其他类型：保持内容不变即可

返回完整的JSON对象。"""

# ════════════════════════════════════════════════════
# 内容生成入口
# ════════════════════════════════════════════════════
def generate_outline(topic, n_slides=13):
    print("  → 生成课程大纲...")
    time.sleep(2)
    user_msg = OUTLINE_USER.format(topic=topic, n_slides=n_slides)
    result = chat([
        {"role": "system", "content": OUTLINE_SYSTEM},
        {"role": "user",   "content": user_msg}
    ])
    return parse_json(result)

def enrich_slide(topic, item):
    """对需要扩充的页面类型进行第二轮增强"""
    enrichable = {"content", "example", "table", "section"}
    if item.get("type") not in enrichable:
        return item
    user_msg = ENRICH_USER.format(topic=topic, outline_item=json.dumps(item, ensure_ascii=False, indent=2))
    try:
        result = chat([
            {"role": "system", "content": ENRICH_SYSTEM},
            {"role": "user",   "content": user_msg}
        ])
        enriched = parse_json(result)
        # 合并，原有字段优先保留结构
        for k, v in enriched.items():
            if v:  # 只更新非空字段
                item[k] = v
    except Exception as e:
        print(f"    ⚠ 扩充失败（{item.get('title','')}）: {e}")
    return item

# ════════════════════════════════════════════════════
# PPT 渲染
# ════════════════════════════════════════════════════
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill, line_color=None, lw=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        if lw: s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=28, bold=False, color=None,
        align=PP_ALIGN.LEFT, wrap=True, italic=False, italic_=False):
    shape = slide.shapes.add_textbox(l, t, w, h)
    tf = shape.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = str(text)
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic or italic_
    run.font.color.rgb = color or C["dark"]
    return shape

def multiline_txt(slide, lines, l, t, w, h, size=28, bold=False,
                  color=None, bullet="", line_spacing=1.3):
    """多行文本，支持 bullet 前缀"""
    shape = slide.shapes.add_textbox(l, t, w, h)
    tf = shape.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = (bullet + line) if bullet else line
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color or C["dark"]
    return shape

# ── 封面 ──────────────────────────────────────────
def render_cover(prs, data, topic):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    # 左深右浅两块
    box(sl, 0, 0, int(W*0.5), H, C["dark"])
    box(sl, int(W*0.5), 0, int(W*0.5), H, C["lgray"])
    # 左底蓝条
    box(sl, 0, H-Inches(0.45), int(W*0.5), Inches(0.45), C["accent"])
    # 右顶蓝条
    box(sl, int(W*0.5), 0, int(W*0.5), Inches(0.45), C["accent"])
    # 主标题
    title = data.get("title", topic)
    txt(sl, title, Inches(1.5), Inches(3.0), int(W*0.44), Inches(3.5),
        size=64, bold=True, color=C["white"])
    # 副标题
    sub = data.get("subtitle","")
    if sub:
        txt(sl, sub, Inches(1.5), Inches(7.0), int(W*0.44), Inches(1.5),
            size=34, color=RGBColor(0xAA,0xBB,0xCC))
    # 左下标注
    txt(sl, "内部培训资料", Inches(1.5), H-Inches(1.5), Inches(10), Inches(0.8),
        size=22, color=C["gray"])
    # 右侧装饰大字
    txt(sl, "TRAINING\nCOURSE", int(W*0.55), Inches(4.0), int(W*0.38), Inches(5.0),
        size=72, bold=True, color=C["line"], italic_=True)

# ── 章节页 ────────────────────────────────────────
def render_section(prs, data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    set_bg(sl, C["white"])
    # 顶部蓝条
    box(sl, 0, 0, W, Inches(0.5), C["accent"])
    # 左侧蓝竖条
    box(sl, 0, Inches(0.5), Inches(0.5), H-Inches(0.5), C["accent"])
    # 右侧内容区浅灰
    box(sl, Inches(0.5), Inches(0.5), W-Inches(0.5), H-Inches(0.5), C["lgray"])

    txt(sl, "— C H A P T E R —",
        Inches(2.0), Inches(2.5), Inches(14), Inches(0.9),
        size=20, color=C["accent"], italic_=True)
    txt(sl, data["title"],
        Inches(2.0), Inches(3.4), int(W*0.4), Inches(4.5),
        size=52, bold=True, color=C["dark"])

    pts = data.get("key_points", []) or data.get("points", [])
    if pts:
        txt(sl, "本章内容",
            int(W*0.5), Inches(2.5), Inches(14), Inches(0.8),
            size=22, bold=True, color=C["gray"])
        for i, pt in enumerate(pts[:5]):
            y = Inches(3.5) + i * Inches(1.65)
            dot = sl.shapes.add_shape(9, int(W*0.5), y+Pt(4), Inches(0.55), Inches(0.55))
            dot.fill.solid(); dot.fill.fore_color.rgb = C["accent"]
            dot.line.fill.background()
            txt(sl, str(i+1), int(W*0.5), y+Pt(4), Inches(0.55), Inches(0.55),
                size=16, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
            txt(sl, pt, int(W*0.5)+Inches(0.75), y, int(W*0.44), Inches(1.3),
                size=28, color=C["dark"])

# ── 内容页 ────────────────────────────────────────
def render_content(prs, data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C["white"])
    W, H = prs.slide_width, prs.slide_height
    # 顶部标题条
    box(sl, 0, 0, W, Inches(1.9), C["accent"])
    box(sl, Inches(0.6), Inches(0.35), Inches(0.1), Inches(1.2), C["white"])
    txt(sl, data["title"], Inches(0.9), Inches(0.38), W-Inches(1.5), Inches(1.15),
        size=40, bold=True, color=C["white"])

    pts = data.get("key_points", []) or data.get("points", [])
    detail = data.get("detail", "")
    n = len(pts)

    if detail and n <= 4:
        lw = int(W*0.55); rx = int(W*0.58); rw = int(W*0.39)
        row_h = min(Inches(1.7), (H-Inches(2.4)) / max(n,1))
        for i, pt in enumerate(pts):
            y = Inches(2.1) + i*(row_h + Inches(0.08))
            # 序号方块
            nb = sl.shapes.add_shape(1, Inches(0.5), y+Pt(4), Inches(1.0), Inches(1.0))
            nb.fill.solid(); nb.fill.fore_color.rgb = C["accent"]
            nb.line.fill.background()
            txt(sl, str(i+1), Inches(0.5), y+Pt(4), Inches(1.0), Inches(1.0),
                size=26, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
            txt(sl, pt, Inches(1.75), y+Pt(8), lw-Inches(2.1), row_h-Pt(12),
                size=29, bold=True, color=C["dark"])
            # 淡灰分隔线
            box(sl, Inches(1.75), y+row_h-Pt(2), lw-Inches(2.2), Pt(2), C["line"])

        # 右侧 detail 卡片
        box(sl, rx, Inches(2.0), rw, H-Inches(2.3), C["lgray"], C["line"], 0.5)
        box(sl, rx, Inches(2.0), Inches(0.14), H-Inches(2.3), C["accent"])
        txt(sl, "📌  详细说明",
            rx+Inches(0.35), Inches(2.2), rw-Inches(0.55), Inches(0.75),
            size=22, bold=True, color=C["accent"])
        txt(sl, detail,
            rx+Inches(0.35), Inches(3.05), rw-Inches(0.55), H-Inches(3.6),
            size=24, color=C["gray"], wrap=True)
    else:
        row_h = min(Inches(1.6), (H-Inches(2.5)) / max(n,1))
        for i, pt in enumerate(pts):
            y = Inches(2.1) + i*(row_h+Inches(0.08))
            bg = C["lgray"] if i%2==0 else C["white"]
            box(sl, Inches(0.4), y, W-Inches(0.8), row_h, bg, C["line"], 0.3)
            box(sl, Inches(0.4), y, Inches(0.14), row_h, C["accent"])
            txt(sl, pt, Inches(0.8), y+Pt(8), W-Inches(1.4), row_h-Pt(16),
                size=29, bold=True, color=C["dark"])
        if detail:
            box(sl, Inches(0.4), H-Inches(2.4), W-Inches(0.8), Inches(2.1),
                C["hilite"], C["line"], 0.5)
            txt(sl, detail, Inches(0.8), H-Inches(2.3), W-Inches(1.6), Inches(1.9),
                size=24, color=C["gray"])

# ── 案例页 ────────────────────────────────────────
def render_example(prs, data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C["white"])
    W, H = prs.slide_width, prs.slide_height
    box(sl, 0, 0, W, Inches(1.9), C["dark"])
    box(sl, Inches(0.6), Inches(0.35), Inches(0.1), Inches(1.2), C["accent2"])
    txt(sl, "📋  " + data["title"], Inches(0.9), Inches(0.38), W-Inches(1.5), Inches(1.15),
        size=40, bold=True, color=C["white"])

    scenario = data.get("scenario","")
    pts = data.get("key_points",[]) or data.get("points",[])

    if scenario:
        bw = int(W*0.57); rx = int(W*0.61); rw = int(W*0.36)
        box(sl, Inches(0.4), Inches(2.0), bw, H-Inches(2.3), C["lgray"], C["line"], 0.5)
        box(sl, Inches(0.4), Inches(2.0), Inches(0.14), H-Inches(2.3), C["accent2"])
        txt(sl, "💬  实战场景",
            Inches(0.75), Inches(2.2), bw-Inches(0.5), Inches(0.75),
            size=22, bold=True, color=C["accent2"])
        txt(sl, scenario, Inches(0.75), Inches(3.1), bw-Inches(0.6), H-Inches(3.7),
            size=25, color=C["dark"], wrap=True)
        if pts:
            txt(sl, "✅  关键要点",
                rx, Inches(2.0), rw, Inches(0.75),
                size=22, bold=True, color=C["accent2"])
            for i, pt in enumerate(pts[:5]):
                y = Inches(3.0) + i*Inches(1.9)
                box(sl, rx, y, rw, Inches(1.7), C["lgray"], C["line"], 0.3)
                box(sl, rx, y, Inches(0.12), Inches(1.7), C["accent2"])
                txt(sl, pt, rx+Inches(0.3), y+Pt(10), rw-Inches(0.45), Inches(1.5),
                    size=25, color=C["dark"])
    else:
        rh = min(Inches(1.6), (H-Inches(2.5))/max(len(pts),1))
        for i, pt in enumerate(pts[:6]):
            y = Inches(2.1) + i*(rh+Inches(0.1))
            box(sl, Inches(0.4), y, W-Inches(0.8), rh, C["lgray"], C["line"], 0.3)
            box(sl, Inches(0.4), y, Inches(0.14), rh, C["accent2"])
            txt(sl, pt, Inches(0.8), y+Pt(8), W-Inches(1.4), rh-Pt(16),
                size=29, bold=True)

# ── 表格页 ────────────────────────────────────────
def render_table(prs, data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C["white"])
    W, H = prs.slide_width, prs.slide_height
    box(sl, 0, 0, W, Inches(1.9), C["dark"])
    box(sl, Inches(0.6), Inches(0.35), Inches(0.1), Inches(1.2), C["accent"])
    txt(sl, data["title"], Inches(0.9), Inches(0.38), W-Inches(1.5), Inches(1.15),
        size=40, bold=True, color=C["white"])

    headers = data.get("table_headers",[])
    rows = data.get("table_rows",[])
    if not headers:
        render_content(prs, data); return

    nc = len(headers); nr = len(rows)
    tt = Inches(2.1); avh = H-Inches(2.4)
    rh = avh / (nr+1); cw = (W-Inches(1.0)) / nc

    for j, h in enumerate(headers):
        x = Inches(0.5) + j*cw
        box(sl, x, tt, cw-Pt(3), rh, C["accent"])
        txt(sl, h, x+Inches(0.12), tt+Pt(6), cw-Inches(0.24), rh-Pt(12),
            size=26, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
    for i, row in enumerate(rows):
        y = tt + (i+1)*rh
        bg = C["lgray"] if i%2==0 else C["white"]
        for j, cell in enumerate(row[:nc]):
            x = Inches(0.5) + j*cw
            box(sl, x, y, cw-Pt(3), rh, bg, C["line"], 0.3)
            txt(sl, str(cell), x+Inches(0.12), y+Pt(4), cw-Inches(0.24), rh-Pt(8),
                size=22, color=C["dark"], align=PP_ALIGN.CENTER)

# ── 问答页 ────────────────────────────────────────
def render_quiz(prs, data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C["white"])
    W, H = prs.slide_width, prs.slide_height
    box(sl, 0, 0, W, Inches(1.9), C["accent"])
    box(sl, Inches(0.6), Inches(0.35), Inches(0.1), Inches(1.2), C["white"])
    txt(sl, "❓  互动问答", Inches(0.9), Inches(0.38), W-Inches(1.5), Inches(1.15),
        size=40, bold=True, color=C["white"])

    q = data.get("quiz_question", data.get("title",""))
    box(sl, Inches(0.4), Inches(2.1), W-Inches(0.8), Inches(1.9),
        C["lgray"], C["line"], 0.5)
    txt(sl, q, Inches(0.8), Inches(2.2), W-Inches(1.6), Inches(1.7),
        size=32, bold=True, align=PP_ALIGN.CENTER)

    opts = data.get("quiz_options",[]) or data.get("points",[])
    ans = data.get("quiz_answer","A")
    ow = (W-Inches(1.5))/2; oh = Inches(2.1)
    for i, opt in enumerate(opts[:4]):
        x = Inches(0.4) + (i%2)*(ow+Inches(0.7))
        y = Inches(4.3) + (i//2)*(oh+Inches(0.45))
        is_a = opt.startswith(ans+".")
        bg = C["green_bg"] if is_a else C["lgray"]
        lc = C["accent2"] if is_a else C["line"]
        box(sl, x, y, ow, oh, bg, lc, 2.0 if is_a else 0.5)
        txt(sl, opt, x+Inches(0.25), y+Pt(12), ow-Inches(0.5), oh-Pt(24),
            size=26, color=C["dark"])

    exp = data.get("quiz_explain","")
    if exp:
        box(sl, Inches(0.4), H-Inches(1.55), W-Inches(0.8), Inches(1.2), C["green_bg"])
        txt(sl, f"✅  正确答案：{ans}  ｜  {exp}",
            Inches(0.8), H-Inches(1.45), W-Inches(1.6), Inches(1.0),
            size=22, color=RGBColor(0x1B,0x5E,0x20))

# ── 总结页 ────────────────────────────────────────
def render_summary(prs, data):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C["white"])
    W, H = prs.slide_width, prs.slide_height
    box(sl, 0, 0, W, Inches(0.45), C["accent"])
    box(sl, 0, H-Inches(0.45), W, Inches(0.45), C["accent"])
    txt(sl, data["title"], Inches(0.5), Inches(0.55), W-Inches(1), Inches(1.8),
        size=46, bold=True, align=PP_ALIGN.CENTER)
    box(sl, W/2-Inches(5), Inches(2.5), Inches(10), Pt(4), C["accent"])

    pts = data.get("key_points",[]) or data.get("points",[])
    n = max(len(pts),1)
    cw = (W-Inches(2.0))/n; ch = Inches(6.2); ct = Inches(2.9)
    for i, pt in enumerate(pts):
        x = Inches(1.0) + i*cw
        box(sl, x+Inches(0.1), ct, cw-Inches(0.2), ch, C["lgray"], C["line"], 0.5)
        box(sl, x+Inches(0.1), ct, cw-Inches(0.2), Inches(1.1), C["accent"])
        txt(sl, str(i+1), x+Inches(0.1), ct, cw-Inches(0.2), Inches(1.1),
            size=32, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        txt(sl, pt, x+Inches(0.2), ct+Inches(1.3), cw-Inches(0.4), ch-Inches(1.5),
            size=24, align=PP_ALIGN.CENTER, wrap=True)

# ════════════════════════════════════════════════════
# 渲染分发
# ════════════════════════════════════════════════════
RENDERERS = {
    "cover":   render_cover,
    "section": render_section,
    "content": render_content,
    "table":   render_table,
    "example": render_example,
    "quiz":    render_quiz,
    "summary": render_summary,
}

def main():
    parser = argparse.ArgumentParser(description="生成培训PPT（两阶段生成）")
    parser.add_argument("--topic",  required=True, help="培训主题")
    parser.add_argument("--output", required=True, help="输出文件 (.pptx)")
    parser.add_argument("--slides", type=int, default=13, help="目标页数（默认13）")
    parser.add_argument("--no-enrich", action="store_true", help="跳过内容扩充阶段")
    args = parser.parse_args()

    if not API_KEY:
        print("ERROR: AIGC_API_KEY 未设置"); sys.exit(1)

    print(f"[1/3] 生成课程大纲（目标 {args.slides} 页）...")
    outline = generate_outline(args.topic, args.slides)
    print(f"      ✓ {len(outline)} 页大纲")

    if not args.no_enrich:
        print(f"[2/3] 扩充各页内容（共 {len(outline)} 页）...")
        enriched = []
        enrich_types = {"content", "example", "table", "section"}
        for i, item in enumerate(outline):
            stype = item.get("type","content")
            if stype in enrich_types:
                print(f"      → [{i+1}/{len(outline)}] {item.get('title','')} ({stype})")
                time.sleep(2)
                item = enrich_slide(args.topic, item)
            enriched.append(item)
        slides_data = enriched
    else:
        slides_data = outline
        print("[2/3] 跳过内容扩充")

    print(f"[3/3] 渲染 PPT（{len(slides_data)} 页）...")
    prs = Presentation()
    prs.slide_width  = Inches(26.6667)
    prs.slide_height = Inches(15.0)

    for s in slides_data:
        stype = s.get("type","content")
        renderer = RENDERERS.get(stype, render_content)
        if stype == "cover":
            renderer(prs, s, args.topic)
        else:
            renderer(prs, s)

    prs.save(args.output)
    print(f"\n✅ 完成！{len(slides_data)} 页 → {args.output}")

if __name__ == "__main__":
    main()
